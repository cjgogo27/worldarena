#!/usr/bin/env python3

from __future__ import annotations

import csv
import json
import math
from pathlib import Path

import imageio.v2 as imageio
import matplotlib.pyplot as plt
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("/data/alice/cjtest/worldmodelphy")
DOCS = BASE / "docs"
RUNS = BASE / "artifacts" / "runs"
FIGS = DOCS / "assets" / "figures"
VIDEOS = DOCS / "assets" / "videos"
TABLE = DOCS / "assets" / "tables" / "table_quantitative_results.csv"
GEN = DOCS / "ppt_generated"
OUT = DOCS / "presentation_worldmodelphy_cn.pptx"

SW = Inches(13.333)
SH = Inches(7.5)

WHITE = RGBColor(251, 249, 246)
PAPER = RGBColor(245, 241, 236)
RED = RGBColor(145, 25, 42)
RED_DARK = RGBColor(112, 18, 32)
INK = RGBColor(25, 25, 28)
MUTED = RGBColor(92, 92, 96)
LINE = RGBColor(223, 218, 211)
SOFT_RED = RGBColor(248, 238, 240)
SOFT_BLUE = RGBColor(238, 242, 247)
SOFT_GOLD = RGBColor(248, 244, 233)
GREEN = RGBColor(34, 110, 77)


def load_table() -> dict[str, dict[str, float | str]]:
    out: dict[str, dict[str, float | str]] = {}
    with TABLE.open(encoding="utf-8") as f:
        for row in csv.DictReader(f):
            parsed: dict[str, float | str] = {"experiment": row["experiment"]}
            for key, value in row.items():
                if key == "experiment":
                    continue
                parsed[key] = float(value) if value.lower() != "nan" else math.nan
            out[str(parsed["experiment"])] = parsed
    return out


def gif_to_mp4(gif_path: Path, fps: int = 4) -> Path:
    GEN.mkdir(parents=True, exist_ok=True)
    mp4 = GEN / f"{gif_path.stem}.mp4"
    if mp4.exists() and mp4.stat().st_mtime >= gif_path.stat().st_mtime:
        return mp4
    frames = imageio.mimread(gif_path)
    writer = imageio.get_writer(mp4, fps=fps, codec="libx264", quality=7, macro_block_size=1)
    for frame in frames:
        writer.append_data(frame)
    writer.close()
    return mp4


def gif_strip(gif_path: Path, picks: int = 4, size: tuple[int, int] = (220, 140)) -> Path:
    GEN.mkdir(parents=True, exist_ok=True)
    out = GEN / f"{gif_path.stem}_strip.png"
    if out.exists() and out.stat().st_mtime >= gif_path.stat().st_mtime:
        return out
    with Image.open(gif_path) as gif:
        total = max(1, getattr(gif, "n_frames", 1))
        ids = sorted({min(total - 1, round(i * (total - 1) / max(picks - 1, 1))) for i in range(picks)})
        imgs: list[Image.Image] = []
        for idx in ids:
            gif.seek(idx)
            frame = gif.convert("RGB")
            frame = ImageOps.contain(frame, size)
            canvas = Image.new("RGB", size, (255, 255, 255))
            canvas.paste(frame, ((size[0] - frame.width) // 2, (size[1] - frame.height) // 2))
            imgs.append(canvas)
    gap = 8
    strip = Image.new("RGB", (len(imgs) * size[0] + (len(imgs) - 1) * gap, size[1]), (255, 255, 255))
    x = 0
    for img in imgs:
        strip.paste(img, (x, 0))
        x += size[0] + gap
    strip.save(out)
    return out


def summary_charts(data: dict[str, dict[str, float | str]]) -> dict[str, Path]:
    GEN.mkdir(parents=True, exist_ok=True)
    motions = ["circular", "projectile", "bounce", "pendulum", "two_body"]
    context = GEN / "chart_context.png"
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.7), dpi=220)
    x = list(range(len(motions)))
    w = 0.36
    short_id = [float(data[f"{m}_gru_short"]["trajectory_mse"]) for m in motions]
    long_id = [float(data[f"{m}_gru_long"]["trajectory_mse"]) for m in motions]
    short_ood = [float(data[f"{m}_gru_short"]["ood_trajectory_mse"]) for m in motions]
    long_ood = [0.0 if math.isnan(float(data[f"{m}_gru_long"]["ood_trajectory_mse"])) else float(data[f"{m}_gru_long"]["ood_trajectory_mse"]) for m in motions]
    axes[0].bar([i - w / 2 for i in x], short_id, width=w, color="#c6a3aa", label="short")
    axes[0].bar([i + w / 2 for i in x], long_id, width=w, color="#91192A", label="long")
    axes[0].set_title("ID trajectory MSE")
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(motions, rotation=20)
    axes[0].legend(frameon=False)
    axes[1].bar([i - w / 2 for i in x], short_ood, width=w, color="#d9c7a7", label="short")
    axes[1].bar([i + w / 2 for i in x], long_ood, width=w, color="#B48C3C", label="long")
    axes[1].set_title("OOD trajectory MSE")
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(motions, rotation=20)
    axes[1].legend(frameon=False)
    fig.tight_layout()
    fig.savefig(context, facecolor="white")
    plt.close(fig)

    arch = GEN / "chart_arch.png"
    fig, ax = plt.subplots(figsize=(11, 4.6), dpi=220)
    subset = ["circular", "projectile", "pendulum"]
    x = list(range(len(subset)))
    w = 0.24
    gru = [float(data[f"{m}_gru_short"]["trajectory_mse"]) for m in subset]
    local = [float(data[f"{m}_local_short"]["trajectory_mse"]) for m in subset]
    bottleneck = [float(data[f"{m}_bottleneck_short"]["trajectory_mse"]) for m in subset]
    ax.bar([i - w for i in x], gru, width=w, color="#91192A", label="GRU")
    ax.bar(x, local, width=w, color="#BE7B63", label="Local attention")
    ax.bar([i + w for i in x], bottleneck, width=w, color="#A99D70", label="Bottleneck")
    ax.set_title("Short-context architecture comparison")
    ax.set_ylabel("ID trajectory MSE")
    ax.set_xticks(x)
    ax.set_xticklabels(subset)
    ax.legend(frameon=False, ncol=3)
    fig.tight_layout()
    fig.savefig(arch, facecolor="white")
    plt.close(fig)

    probe = GEN / "chart_probe_gap.png"
    selected = ["circular_gru_long", "projectile_gru_long", "pendulum_gru_long", "two_body_gru_long", "circular_local_short", "circular_bottleneck_short"]
    fig, ax = plt.subplots(figsize=(11.5, 4.8), dpi=220)
    x = list(range(len(selected)))
    w = 0.35
    idv = [float(data[s]["probe_r2_mean"]) for s in selected]
    oodv = [float(data[s]["ood_probe_r2_mean"]) for s in selected]
    ax.bar([i - w / 2 for i in x], idv, width=w, color="#91192A", label="ID probe")
    ax.bar([i + w / 2 for i in x], oodv, width=w, color="#D4B483", label="OOD probe")
    ax.axhline(0, color="#888", linewidth=1)
    ax.set_title("ID probe can be perfect while OOD probe collapses")
    ax.set_xticks(x)
    ax.set_xticklabels(selected, rotation=18, ha="right")
    ax.legend(frameon=False)
    fig.tight_layout()
    fig.savefig(probe, facecolor="white")
    plt.close(fig)
    return {"context": context, "arch": arch, "probe": probe}


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def add_bg(slide, page: int, section: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SW, Inches(0.08))
    top.fill.solid()
    top.fill.fore_color.rgb = RED
    top.line.fill.background()
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.6), Inches(0.18), Inches(1.1), Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = PAPER
    tag.line.color.rgb = LINE
    tf = tag.text_frame
    tf.paragraphs[0].text = section
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    num = slide.shapes.add_textbox(Inches(12.65), Inches(6.98), Inches(0.3), Inches(0.2))
    p = num.text_frame.paragraphs[0]
    p.text = str(page)
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def title(slide, title_text: str, kicker: str | None = None, subtitle: str | None = None) -> None:
    if kicker:
        tb = slide.shapes.add_textbox(Inches(0.72), Inches(0.28), Inches(4.0), Inches(0.22))
        p = tb.text_frame.paragraphs[0]
        p.text = kicker
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RED
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(0.5), Inches(11.4), Inches(0.72))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = INK
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.14), Inches(1.15), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(11.2), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED


def card(slide, x, y, w, h, header: str, lines: list[str], fill=WHITE, border=LINE, header_fill=None, header_color=INK, body_size=13):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = fill
    outer.line.color.rgb = border
    outer.line.width = Pt(1.1)
    hf = header_fill if header_fill is not None else fill
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.42))
    band.fill.solid()
    band.fill.fore_color.rgb = hf
    band.line.fill.background()
    ht = slide.shapes.add_textbox(x + Inches(0.24), y + Inches(0.16), w - Inches(0.48), Inches(0.24))
    hp = ht.text_frame.paragraphs[0]
    hp.text = header
    hp.font.size = Pt(15)
    hp.font.bold = True
    hp.font.color.rgb = header_color
    tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.68), w - Inches(0.44), h - Inches(0.82))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line.startswith("•") else f"• {line}"
        p.font.size = Pt(body_size)
        p.font.color.rgb = INK
        p.space_after = Pt(4)


def picture_card(slide, image: Path, x, y, w, h, caption: str | None = None):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(1.0)
    pic_h = h - Inches(0.26 if caption else 0.16)
    slide.shapes.add_picture(str(image), x + Inches(0.08), y + Inches(0.08), width=w - Inches(0.16), height=pic_h)
    if caption:
        tb = slide.shapes.add_textbox(x + Inches(0.14), y + h - Inches(0.22), w - Inches(0.28), Inches(0.16))
        p = tb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(9.5)
        p.font.color.rgb = MUTED


def movie_card(slide, movie: Path, poster: Path, x, y, w, h, caption: str | None = None):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE
    box.line.width = Pt(1.0)
    slide.shapes.add_movie(str(movie), x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), h - Inches(0.26 if caption else 0.16), poster_frame_image=str(poster), mime_type="video/mp4")
    if caption:
        tb = slide.shapes.add_textbox(x + Inches(0.14), y + h - Inches(0.22), w - Inches(0.28), Inches(0.16))
        p = tb.text_frame.paragraphs[0]
        p.text = caption
        p.font.size = Pt(9.5)
        p.font.color.rgb = MUTED


def footer(slide, text: str):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(6.74), Inches(11.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.83), Inches(11.2), Inches(0.26))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RED


def run_metrics(name: str) -> dict:
    return json.loads((RUNS / name / "metrics.json").read_text())


def fmt(x: float | int) -> str:
    if isinstance(x, float) and math.isnan(x):
        return "NaN"
    return f"{x:.2f}"


def run_slide(prs: Presentation, page: int, section: str, run_name: str, claim: str, takeaway: str, bullets: list[str]) -> None:
    metrics = run_metrics(run_name)
    run_dir = RUNS / run_name
    test_gif = run_dir / "test_rollout.gif"
    ood_gif = run_dir / "ood_rollout.gif"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, section)
    title(slide, claim, run_name.replace("_", " / ").upper(), f"ID trajectory={fmt(metrics['test_metrics']['trajectory_mse'])} | OOD trajectory={fmt(metrics['ood_metrics']['trajectory_mse'])} | ID probe={fmt(metrics['probe_test']['probe_r2_mean'])} | OOD probe={fmt(metrics['probe_ood']['probe_r2_mean'])}")
    movie_card(slide, gif_to_mp4(test_gif), gif_strip(test_gif), Inches(0.78), Inches(1.72), Inches(5.7), Inches(2.15), "测试集 rollout（可播放视频）")
    movie_card(slide, gif_to_mp4(ood_gif), gif_strip(ood_gif), Inches(6.82), Inches(1.72), Inches(5.7), Inches(2.15), "OOD rollout（可播放视频）")
    picture_card(slide, run_dir / "loss_curve.png", Inches(0.78), Inches(4.18), Inches(4.25), Inches(2.18), "训练/验证 loss")
    picture_card(slide, run_dir / "ood_rollout_strip.png", Inches(5.22), Inches(4.18), Inches(3.05), Inches(2.18), "OOD 关键帧条带")
    card(slide, Inches(8.48), Inches(4.18), Inches(4.04), Inches(2.18), "这页想说明什么", bullets, fill=WHITE, border=LINE, header_fill=SOFT_RED, header_color=RED, body_size=12.5)
    footer(slide, takeaway)


def build() -> Path:
    data = load_table()
    charts = summary_charts(data)
    prs = prs_new()
    page = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Intro")
    title(slide, "从开普勒到牛顿：视频生成模型内部是否涌现了物理？", "WORLDMODELPHY", "基于合成视频、运行级别 artifact、OOD 评测与文献证据链的答辩/组会版汇报")
    movie_card(slide, gif_to_mp4(RUNS / "circular_gru_long" / "test_rollout.gif"), gif_strip(RUNS / "circular_gru_long" / "test_rollout.gif"), Inches(7.12), Inches(1.28), Inches(5.38), Inches(2.15), "ID：circular / GRU-long")
    movie_card(slide, gif_to_mp4(RUNS / "projectile_gru_long" / "ood_rollout.gif"), gif_strip(RUNS / "projectile_gru_long" / "ood_rollout.gif"), Inches(7.12), Inches(3.76), Inches(5.38), Inches(2.15), "OOD：projectile / GRU-long")
    card(slide, Inches(0.78), Inches(4.95), Inches(5.9), Inches(0.95), "先给结论", ["弱形式“物理感”存在：模型可以学到 motion statistics 和部分局部动力学 cue", "强形式 physics world model 证据不足：一旦 OOD，rollout 和 probe 都会迅速崩坏"], fill=SOFT_RED, border=RED, header_fill=SOFT_RED, header_color=RED)
    footer(slide, "这次汇报的重点不是“视频能不能生成”，而是“模型内部到底学成了什么”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Story")
    title(slide, "我们要按三个层次回答：理论启发、领域现状、以及我们自己的像素空间实验", "AGENDA")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "问题 1：From Kepler to Newton", ["论文到底证明了什么", "为什么 temporal locality 是关键", "curve-fitting 与 learning dynamics 的边界在哪里"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "问题 2：当前 Video Generative Model", ["哪些 benchmark 最能回答“是否涌现了物理”", "它们到底测了什么", "主流模型为何仍然普遍不过关"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "问题 3：像素输入的小模型", ["模型可以多小，为什么不能只谈参数量", "long context 为什么会带来高 ID / 低 OOD", "如何用 GIF、图和 probe 一起判断物理是否真的出现"], header_fill=SOFT_GOLD)
    footer(slide, "整体结构采用 onion model：why → evidence → implication。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q1")
    title(slide, "`From Kepler to Newton` 的核心不是“预测得准”，而是“模型内部到底学成了曲线还是动力学”", "QUESTION 1")
    card(slide, Inches(0.78), Inches(1.72), Inches(5.72), Inches(4.05), "Keplerian 表征", ["更像记住全局轨道几何 / 曲线形状", "可以在训练分布内拿到很强的预测结果", "但并不要求模型掌握局部动力学机制"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(6.82), Inches(1.72), Inches(5.72), Inches(4.05), "Newtonian 表征", ["更像编码局部力、速度、加速度关系", "依赖当前局部状态推出未来", "更接近可迁移、可干预的 world model"], header_fill=SOFT_BLUE)
    card(slide, Inches(0.78), Inches(6.0), Inches(11.76), Inches(0.58), "对我们最重要的启发", ["如果不限制模型背整段轨迹的自由度，它很容易以高精度假装“懂物理”；因此 short vs long、OOD、probe 都是必要的。"], fill=SOFT_GOLD, border=RGBColor(218, 198, 156), header_fill=SOFT_GOLD)
    footer(slide, "这也是我们后面所有设计的理论起点。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q2")
    title(slide, "文献已经给出很一致的结论：主流视频模型“像真的”，但普遍还不“按物理规律地真”", "QUESTION 2")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(2.0), "Physics-IQ · ICCV 2025", ["做了什么：396 个真实视频，覆盖流体/固体/光学/热学/磁学", "怎么测：只有真正理解物理才可能完成的视频生成任务", "结论：visual realism 与 physics understanding 弱相关"], header_fill=SOFT_RED, header_color=RED, body_size=12.4)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(2.0), "Morpheus · 2025", ["做了什么：130 个真实物理实验视频", "怎么测：conservation-law / physics-informed metrics", "结论：生成视频很美，但经常违背物理约束"], header_fill=SOFT_RED, header_color=RED, body_size=12.4)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(2.0), "WorldScore · ICCV 2025", ["做了什么：统一评测 quality / controllability / dynamics", "怎么测：3000 个测试案例", "结论：质量强，但 dynamics 与 control 不足"], header_fill=SOFT_RED, header_color=RED, body_size=12.4)
    card(slide, Inches(0.78), Inches(4.1), Inches(5.78), Inches(1.86), "PhyGenBench · ICML 2025", ["做了什么：160 个 prompts、27 条物理常识 / 规律", "结论：当前 T2V 模型普遍难以稳定生成符合物理常识的视频"], header_fill=SOFT_BLUE, body_size=12.4)
    card(slide, Inches(6.76), Inches(4.1), Inches(5.78), Inches(1.86), "WorldBench / PhyWorldBench · 2025–2026", ["做了什么：把 physical world evaluation 拆成 motion / permanence / support / scale 等子概念", "结论：最佳模型仍离“可靠物理世界模型”有明显差距"], header_fill=SOFT_GOLD, body_size=12.4)
    footer(slide, "这些 benchmark 的一致信息是：会生成“像真的”视频，不等于会“按物理规律泛化”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q2")
    title(slide, "更关键的是，几篇论文已经明确告诉我们：模型为什么会“假装懂物理”", "WHY CURRENT MODELS CHEAT")
    card(slide, Inches(0.78), Inches(1.72), Inches(5.82), Inches(2.2), "From Kepler to Newton · 2026", ["怎么做：系统引入 spatial smoothness、stability、temporal locality 三类最小偏置", "怎么测：模型到底学成 Keplerian curve-fitting，还是 Newtonian force representation", "结论：没有 temporal locality，Transformer 很容易靠背曲线拿高分"], header_fill=SOFT_RED, header_color=RED, body_size=12.4)
    card(slide, Inches(6.72), Inches(1.72), Inches(5.82), Inches(2.2), "PISA / PisaBench · ICML 2025", ["怎么做：对自由落体等简单牛顿运动做 physics post-training", "怎么测：trajectory L2、collision consistency 等指标", "结论：后训练能修一点，但 generalization 仍然没有解决"], header_fill=SOFT_BLUE, body_size=12.4)
    card(slide, Inches(0.78), Inches(4.25), Inches(5.82), Inches(1.75), "PhyT2V · CVPR 2025", ["怎么做：用 LLM CoT + iterative refinement 辅助视频生成", "结论：外部 reasoning 能改善物理遵从性，但这不是内部自然涌现的证据"], header_fill=SOFT_GOLD, body_size=12.4)
    card(slide, Inches(6.72), Inches(4.25), Inches(5.82), Inches(1.75), "这对我们自己的实验意味着什么", ["必须把“像不像视频”与“是不是物理上对”拆开", "必须看 short/long、OOD、probe 与失败案例", "这就是为什么后面每页都用运行级 artifact 说话"], fill=SOFT_RED, border=RED, header_fill=SOFT_RED, header_color=RED_DARK, body_size=12.4)
    footer(slide, "我们的小模型实验不是重复 benchmark，而是在可控像素空间里拆开失败原因。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q3")
    title(slide, "我们把坐标输入换成图片后，问题立刻变成：模型必须同时学感知与动力学", "QUESTION 3")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.72), "实验矩阵", ["分辨率 64×64，帧长 32，灰度纯背景", "任务：circular / projectile / bounce / pendulum / two_body", "模型：GRU / local-attention / bottleneck", "train / val / test / OOD = 96 / 24 / 24 / 24"], header_fill=SOFT_BLUE)
    picture_card(slide, gif_strip(RUNS / "circular_gru_long" / "test_rollout.gif"), Inches(4.82), Inches(1.72), Inches(2.35), Inches(0.9), "circular")
    picture_card(slide, gif_strip(RUNS / "projectile_gru_long" / "test_rollout.gif"), Inches(7.36), Inches(1.72), Inches(2.35), Inches(0.9), "projectile")
    picture_card(slide, gif_strip(RUNS / "bounce_gru_long" / "test_rollout.gif"), Inches(9.90), Inches(1.72), Inches(2.35), Inches(0.9), "bounce")
    picture_card(slide, gif_strip(RUNS / "pendulum_gru_long" / "test_rollout.gif"), Inches(4.82), Inches(2.9), Inches(3.72), Inches(0.95), "pendulum")
    picture_card(slide, gif_strip(RUNS / "two_body_gru_long" / "test_rollout.gif"), Inches(8.8), Inches(2.9), Inches(3.45), Inches(0.95), "two-body")
    card(slide, Inches(4.82), Inches(4.2), Inches(7.43), Inches(2.24), "这件事为什么比坐标输入难很多", ["模型得先从像素里恢复位置 / 速度 / 相位等隐状态", "再用这些隐状态去外推未来", "因此“能不能看懂画面”与“能不能学会物理”叠在一起"], header_fill=SOFT_GOLD)
    footer(slide, "坐标输入只是在学动力学；像素输入是在学“感知 + 动力学”的组合问题。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q3a")
    title(slide, "模型可以做得很小，但小模型同样会变成“高拟合的轨迹模板机”", "QUESTION 3(a)")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.72), Inches(4.1), "我们实际观察到", ["sub-1M 量级已足够把训练分布拟合得很好", "因此“能不能学会视频”并不是瓶颈", "真正的门槛在于 OOD 与 latent 是否仍保持物理一致"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.78), Inches(1.72), Inches(3.72), Inches(4.1), "为什么不能只谈参数量", ["像素输入增加了感知负担", "模型必须恢复 state，再外推 state", "所以表征容量 + 时序容量共同决定下界"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.78), Inches(1.72), Inches(3.72), Inches(4.1), "真正该问的问题", ["不是“多小能跑通”", "而是“多小还能保持 OOD physics”", "这也是为什么后面每页都强调 ID 与 OOD 的反差"], header_fill=SOFT_GOLD)
    footer(slide, "参数小不代表更接近 physics；小模型同样会走向模板记忆。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q3b")
    title(slide, "GRU baseline 给出最稳定的模式：long context 几乎总能提升 ID，但不能解决 OOD physics", "QUESTION 3(b)")
    picture_card(slide, charts["context"], Inches(0.78), Inches(1.72), Inches(7.0), Inches(4.25), "GRU baseline: short vs long across 5 motions")
    card(slide, Inches(8.05), Inches(1.72), Inches(4.48), Inches(4.25), "从总图先读出三件事", ["Circular / Pendulum / Two-body 的 ID 在 long context 下显著更好", "Bounce 的 OOD 在 long context 下反而更差", "Projectile-long 直接在 OOD collapse，说明高 ID 并不意味着懂动力学"], header_fill=SOFT_RED, header_color=RED)
    footer(slide, "这页是整套实验的总领结论：long context 更像“轨迹模板增强器”。")
    page += 1

    run_slide(prs, page, "Run", "circular_gru_long", "Circular：long context 的确让轨迹更平滑，但 OOD 仍说明它没有学到稳健物理", "Circular 的最佳解释不是“学会了圆周物理”，而是“把训练分布内的圆轨迹拟合得更像了”。", ["ID trajectory MSE 从 83.53 降到 3.24，很强", "但 OOD 仍有 390.34，并没有得到稳健外推", "ID probe=1.00、OOD probe≈0.02，进一步说明 latent 只在训练流形内可解释"]) 
    page += 1
    run_slide(prs, page, "Run", "projectile_gru_long", "Projectile：这是最直观的反例——分布内非常强，不代表分布外还“懂重力”", "Projectile-long 是最应该在答辩时重点讲的失败案例：测试集几乎完美，但 OOD rollout 直接崩。", ["ID trajectory MSE=5.54，看起来像“学会了抛体运动”", "但 OOD trajectory MSE=NaN，说明一旦参数变了就失去稳定演化", "ID probe=1.00 而 OOD probe=-24.91，证明 latent 没有真正学会可迁移动力学"]) 
    page += 1
    run_slide(prs, page, "Run", "bounce_gru_long", "Bounce：碰撞把模板学习暴露得更彻底——ID 改善并没有带来更好的边界动力学", "Bounce 说明离散事件（碰撞、反弹）比光滑轨迹更能戳穿“看起来像对的”模型。", ["long 的 ID 从 71.36 改善到 35.69", "但 OOD 从 644.01 恶化到 1303.67", "碰撞类任务最容易把“不是学物理，只是记模板”这件事暴露出来"]) 
    page += 1
    run_slide(prs, page, "Run", "pendulum_gru_long", "Pendulum：周期系统最容易被“背成形状”，所以它非常像 Keplerian curve fitting 的像素版本", "Pendulum-long 的故事是：ID 几乎完美，但它依然没有变成 Newtonian world model。", ["ID trajectory MSE=1.97，非常强", "OOD trajectory MSE 仍高达 703.66", "周期结构特别适合被记成整体 shape，因此最支持“曲线拟合而非动力学”解释"]) 
    page += 1
    run_slide(prs, page, "Run", "two_body_gru_long", "Two-body 是当前最“像成功”的任务，但它更像学会了整体结构，而不是对象级守恒律", "Two-body 最值得谨慎讲：它确实是最稳的一组，但当前评测还没有对象级守恒量，所以结论不能说过头。", ["ID trajectory MSE=0.33，OOD=208.00，整体上比其他任务更稳", "ID / OOD probe 也相对最好", "但当前评测使用的是简化中心量，因此只能说模型学到结构稳定性，不能说学会了完整 two-body physics"]) 
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Models")
    title(slide, "Local attention 和 bottleneck 都没有自动把模型推向“更物理”的方向", "ARCHITECTURE COMPARISON")
    picture_card(slide, charts["arch"], Inches(0.78), Inches(1.72), Inches(6.8), Inches(4.2), "Short-context architecture comparison")
    movie_card(slide, gif_to_mp4(RUNS / "circular_local_short" / "ood_rollout.gif"), gif_strip(RUNS / "circular_local_short" / "ood_rollout.gif"), Inches(7.9), Inches(1.72), Inches(4.62), Inches(1.75), "local / circular / OOD")
    movie_card(slide, gif_to_mp4(RUNS / "circular_bottleneck_short" / "ood_rollout.gif"), gif_strip(RUNS / "circular_bottleneck_short" / "ood_rollout.gif"), Inches(7.9), Inches(3.84), Inches(4.62), Inches(1.75), "bottleneck / circular / OOD")
    footer(slide, "这页想说明：改结构本身不够，真正缺的是更强的 physics inductive bias 与更严的 intervention-style training/eval。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Q3c")
    title(slide, "如果只看 ID probe，会系统性高估“模型学到了物理”", "QUESTION 3(c)")
    picture_card(slide, charts["probe"], Inches(0.78), Inches(1.72), Inches(7.2), Inches(4.25), "ID probe can be perfect while OOD probe collapses")
    card(slide, Inches(8.25), Inches(1.72), Inches(4.28), Inches(4.25), "正确的评测层次", ["第一层：frame / trajectory / velocity / acceleration", "第二层：ID + OOD linear probe", "第三层：intervention（改重力、初速度、restitution、质量比）", "只有三层同时成立，才能更接近说“学到了物理”"], header_fill=SOFT_BLUE)
    footer(slide, "所以“视频好看”或“ID probe 很高”都不够，必须要求 OOD 下也能成立。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Results")
    title(slide, "把最关键的数字并到一页后，答案会变得非常直接", "KEY QUANTITATIVE SUMMARY")
    table = slide.shapes.add_table(11, 5, Inches(0.78), Inches(1.72), Inches(11.75), Inches(4.75)).table
    headers = ["实验", "ID traj", "OOD traj", "ID probe", "OOD probe"]
    for i, h in enumerate(headers):
        c = table.cell(0, i)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = RED
        p = c.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    rows = [
        ["circular_gru_long", "3.24", "390.34", "1.00", "0.02"],
        ["projectile_gru_long", "5.54", "NaN", "1.00", "-24.91"],
        ["bounce_gru_long", "35.69", "1303.67", "1.00", "0.03"],
        ["pendulum_gru_long", "1.97", "703.66", "1.00", "-0.05"],
        ["two_body_gru_long", "0.33", "208.00", "1.00", "0.16"],
        ["circular_local_short", "311.21", "945.99", "1.00", "0.07"],
        ["projectile_local_short", "1964.77", "898.36", "1.00", "-37.62"],
        ["circular_bottleneck_short", "230.12", "1150.23", "0.95", "-0.02"],
        ["projectile_bottleneck_short", "374.57", "506.44", "0.77", "-13.51"],
        ["pendulum_bottleneck_short", "182.09", "973.19", "0.97", "-0.05"],
    ]
    for r, row in enumerate(rows, start=1):
        for cidx, val in enumerate(row):
            c = table.cell(r, cidx)
            c.text = val
            c.fill.solid()
            c.fill.fore_color.rgb = PAPER if r % 2 else WHITE
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = INK
            if cidx == 0:
                p.font.bold = True
    footer(slide, "一眼就能看出：ID improvement 不自动意味着 OOD physics improvement。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Answer")
    title(slide, "现在可以更精确地回答：我们看到的是 physics-flavored representations，而不是 robust physics world model", "FINAL ANSWER")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "我们确实看到了什么", ["模型能学到运动统计规律", "在 ID 上形成可解码的 latent dynamics", "在某些任务上保持很强的视觉结构"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "我们还没有看到什么", ["稳健的 OOD 物理泛化", "对象级守恒律证据", "对干预变量的一致 latent response"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "所以最终判断", ["弱涌现：有", "强涌现：证据不足", "更像 trajectory template + local dynamics cues 的组合"], header_fill=SOFT_GOLD)
    footer(slide, "一句话：会“像物理地运动”，不等于会“以物理规律的方式泛化”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Next")
    title(slide, "下一步最值得做的不是更大模型，而是更强的 physics inductive bias 与 intervention-based evaluation", "WHAT TO DO NEXT")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.65), "建模", ["在 bottleneck 上加入显式 dynamics loss", "把 local bias 做成真正的 state transition prior", "减少单纯像素重建对模板学习的诱导"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.65), "评测", ["升级 two-body 到对象级守恒量", "把 OOD probe 设成必要标准", "系统加入 counterfactual interventions"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.65), "汇报 takeaway", ["Kepler → Newton 的问题在像素空间依然成立", "long context 很容易带来“高 ID、低 OOD”", "这套 deck 的价值就是把失败原因讲清楚"], header_fill=SOFT_GOLD)
    footer(slide, "如果答辩时只能记住一句话：高 ID + 高视觉质量，并不能证明物理已经涌现。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "后续主线：从“模型为什么能泛化”一直走到“怎样得到真正满足物理规律的世界模型”", "CONTINUATION")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "起点", ["如果生成模型本质上只是记忆训练集，后面的“学到物理”讨论都站不住", "所以必须先问：它为什么能生成训练集之外的新样本？"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "中间问题", ["即便它会泛化，也要追问：它泛化的到底是视觉模式、运动模式，还是更深层物理结构？", "这一步决定我们后续该改的是 loss、表示，还是整个生成机制。"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "终点", ["最终目标不是更强视频预测器，而是能按物理原因生成、并稳定 rollout 的可控世界模型。", "latent action、belief update 和更物理的表示，都是为这个目标服务。"], header_fill=SOFT_GOLD)
    footer(slide, "这一部分不是重复前面的实验，而是把实验结论推进成一条完整研究路线。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "第一层问题：生成模型为什么不是纯记忆器？", "GENERALIZATION")
    card(slide, Inches(0.78), Inches(1.72), Inches(5.72), Inches(4.35), "为什么这是第一步", ["如果模型只是记住训练集、在附近做插值，那后面所有“学到物理”的讨论都不成立。", "必须先解释：为什么 diffusion / flow matching 在现实里会生成训练集之外的样本？"], header_fill=SOFT_BLUE)
    card(slide, Inches(6.82), Inches(1.72), Inches(5.72), Inches(4.35), "真正要研究的机制", ["它没有完全 overfit，是因为网络的 spectral bias？", "还是训练中的噪声、优化路径、隐式正则在阻止完美记忆？", "这一步研究的是：模型为什么不是纯记忆器。"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(0.78), Inches(6.2), Inches(11.76), Inches(0.55), "对现有实验的衔接", ["我们已经看到小模型能够在 ID 内形成结构性表征；下一步就是追问这种结构性到底从哪里来。"], fill=SOFT_GOLD, border=RGBColor(218, 198, 156), header_fill=SOFT_GOLD)
    footer(slide, "先搞清楚“为什么能泛化”，后面谈“是否学到物理”才有基础。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "第二层问题：它泛化的到底是什么——纹理、几何、运动，还是物理机制？", "WHAT IS GENERALIZED")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "可能只是视觉层", ["纹理统计", "低频几何形状", "时空平滑性", "这些都足以让视频“看起来像对的”"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "也可能接近结构层", ["object identity", "transport / correspondence", "局部动力学", "长时一致性"], header_fill=SOFT_BLUE)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "关键判别", ["如果泛化的是纹理和低频形状，那只是视觉生成。", "如果泛化的是 object-level motion 与局部动力学，那才真正开始接近物理。"], header_fill=SOFT_GOLD)
    footer(slide, "所以“模型会泛化”还不够，必须继续问“它按什么结构泛化”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "第三层问题：现有 diffusion / flow 的生成原理，可能天然就不够“物理”", "GENERATION PRINCIPLE")
    card(slide, Inches(0.78), Inches(1.72), Inches(5.72), Inches(4.35), "你真正的不满是什么", ["当前很多生成模型更像是在固定空间坐标 (x, y) 上改像素值 z。", "它更像“原地改颜色”，而不是“物体在空间里移动并保持 identity”。"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(6.82), Inches(1.72), Inches(5.72), Inches(4.35), "为什么这与物理错位", ["真实世界里的变化来自实体运动、接触、相互作用和传输。", "如果底层生成机制主要是像素统计重构，那么它天然更容易学到视觉相关性，而不是物理因果。"], header_fill=SOFT_BLUE)
    card(slide, Inches(0.78), Inches(6.2), Inches(11.76), Inches(0.55), "这一步的结论", ["如果现有 diffusion/flow 的泛化基础主要还是像素统计，那它的生成原理就可能天然不够物理。"], fill=SOFT_GOLD, border=RGBColor(218, 198, 156), header_fill=SOFT_GOLD)
    footer(slide, "这一步把问题从“模型为什么失败”推进到“现有生成方式本身是否错位”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "所以更合理的方向，是把生成单元从“像素值”推进到“实体 / 点 / transport / 状态演化”", "MORE PHYSICAL REPRESENTATION")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "候选生成单元", ["点", "粒子", "object token", "slot", "correspondence / transport"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "为什么它们更接近物理", ["连续性更自然", "identity preservation 更自然", "长时一致性更容易出现", "局部相互作用更容易表达"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "这一步真正要得到什么", ["不是更像的像素视频", "而是更像“实体在空间中持续演化”的生成机制", "这样世界模型才更可能学到局部动力学与更深层原因"], header_fill=SOFT_GOLD)
    footer(slide, "从像素生成走向实体/transport 生成，是整条主线里最关键的结构性转向。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "世界模型的目标不应停留在“学会反应”，而应进一步走到“学会原因”", "FROM DYNAMICS TO CAUSES")
    card(slide, Inches(0.78), Inches(1.72), Inches(5.72), Inches(4.35), "前一阶段只能学到什么", ["下一步怎么动", "短时局部动力学", "接触后会发生什么样的表面变化"], header_fill=SOFT_BLUE)
    card(slide, Inches(6.82), Inches(1.72), Inches(5.72), Inches(4.35), "真正还要恢复的隐变量", ["柔性", "材质", "质量", "摩擦", "支撑关系", "接触条件"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(0.78), Inches(6.2), Inches(11.76), Inches(0.55), "这一步的意义", ["模型不仅知道“会发生什么”，还要逐渐知道“为什么会发生”。"], fill=SOFT_GOLD, border=RGBColor(218, 198, 156), header_fill=SOFT_GOLD)
    footer(slide, "这一步把问题从几何/运动规律推进到因果化的物理机制。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "latent action 在这里出现：它不是起点，而是连接长时语义与局部物理的接口", "LATENT ACTION AS INTERFACE")
    card(slide, Inches(0.78), Inches(1.72), Inches(3.75), Inches(4.7), "语义侧", ["抓", "推", "展开", "试探", "协同操作"], header_fill=SOFT_BLUE)
    card(slide, Inches(4.79), Inches(1.72), Inches(3.75), Inches(4.7), "物理侧", ["接触后的局部变化", "柔性物体的短时响应", "局部稳定性与连续性"], header_fill=SOFT_RED, header_color=RED)
    card(slide, Inches(8.80), Inches(1.72), Inches(3.75), Inches(4.7), "为什么它必须可修正", ["初始 latent 只抓到粗粒度语义", "环境反馈会暴露新的物理信息", "belief 更新后，latent 也应 refinement"], header_fill=SOFT_GOLD)
    footer(slide, "latent action 最好的理解不是“压缩动作”，而是“语义—物理之间的中间变量”。")
    page += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, page, "Part II")
    title(slide, "整条路线最终闭合：从“为什么能泛化”走到“怎样得到一个符合物理规律的可控世界模型”", "UNIFIED ROADMAP")
    card(slide, Inches(0.78), Inches(1.72), Inches(11.76), Inches(4.75), "统一主线（一句话版本）", ["先研究生成模型为什么能泛化，证明它不是纯记忆器；", "再研究它泛化的到底是视觉模式还是物理结构；", "如果现有 diffusion/flow 的泛化基础主要还是像素统计，就说明它的底层生成原理可能不够物理；", "于是需要更接近实体运动、transport 和状态演化的生成方式；", "在这种表示下，世界模型才更可能学到局部动力学与更深层的物理原因；", "为了把这些物理规律真正用于决策和控制，引入 latent action 作为语义—物理接口，并在 belief update 的闭环中不断 refinement；", "最终得到一个能够稳定 rollout、并生成满足物理规律结果的世界模型。"], header_fill=SOFT_GOLD, body_size=13)
    footer(slide, "这就是“后续部分”与前面实验的自然衔接：前面说明现状不够，后面说明下一代路线应该往哪里走。")
    page += 1

    appendix_runs = [
        ("circular_gru_short", "circular_gru_long"),
        ("projectile_gru_short", "projectile_gru_long"),
        ("bounce_gru_short", "bounce_gru_long"),
        ("pendulum_gru_short", "pendulum_gru_long"),
        ("two_body_gru_short", "two_body_gru_long"),
        ("circular_local_short", "circular_bottleneck_short"),
    ]
    for left_run, right_run in appendix_runs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, page, "Appendix")
        title(slide, f"Appendix：{left_run} vs {right_run}", "RUN GALLERY", "备份页：用于答疑时直接播放 run-level 媒体并解释对应指标")
        ldir = RUNS / left_run
        rdir = RUNS / right_run
        movie_card(slide, gif_to_mp4(ldir / "test_rollout.gif"), gif_strip(ldir / "test_rollout.gif"), Inches(0.78), Inches(1.72), Inches(5.75), Inches(1.8), f"{left_run} / test")
        movie_card(slide, gif_to_mp4(rdir / "test_rollout.gif"), gif_strip(rdir / "test_rollout.gif"), Inches(6.78), Inches(1.72), Inches(5.75), Inches(1.8), f"{right_run} / test")
        movie_card(slide, gif_to_mp4(ldir / "ood_rollout.gif"), gif_strip(ldir / "ood_rollout.gif"), Inches(0.78), Inches(3.84), Inches(5.75), Inches(1.8), f"{left_run} / OOD")
        movie_card(slide, gif_to_mp4(rdir / "ood_rollout.gif"), gif_strip(rdir / "ood_rollout.gif"), Inches(6.78), Inches(3.84), Inches(5.75), Inches(1.8), f"{right_run} / OOD")
        footer(slide, "备份页的作用不是堆素材，而是让答疑时能即时切到真正的 run-level 证据。")
        page += 1

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
